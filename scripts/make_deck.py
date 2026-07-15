"""
BPMN IQ 2  –  Executive Presentation
AT&T colour palette  |  python-pptx 1.x
Output: c:\code\BPMN_IQ_2\BPMN_IQ2_Presentation.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import pptx.oxml.ns as nsmap
from lxml import etree
import copy, os

# ── Slide dimensions (16:9 widescreen) ────────────────────────────────────────
W = Inches(13.33)
H = Inches(7.5)

# ── AT&T colour palette ────────────────────────────────────────────────────────
NAVY   = RGBColor(0x00, 0x28, 0x5A)
BLUE   = RGBColor(0x00, 0x9F, 0xDB)
CYAN   = RGBColor(0x00, 0xC3, 0xE3)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LGREY  = RGBColor(0xD5, 0xD6, 0xD2)
MGREY  = RGBColor(0x8E, 0x8E, 0x8E)
DGREY  = RGBColor(0x2C, 0x2C, 0x2C)
ACCENT = RGBColor(0xFF, 0xB8, 0x00)   # amber highlight

ICON_PATH   = r"c:\code\BPMN_IQ_2\client\public\bpmniq_icon.png"
LOGO_PATH   = r"c:\code\BPMN_IQ_2\client\public\bpmniq_logo.png"
OUT_PATH    = r"c:\code\BPMN_IQ_2\BPMN_IQ2_Presentation.pptx"

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

blank_layout = prs.slide_layouts[6]   # completely blank

# ══════════════════════════════════════════════════════════════════════════════
#  LOW-LEVEL HELPERS
# ══════════════════════════════════════════════════════════════════════════════

def add_rect(slide, left, top, width, height, fill_rgb, alpha=None):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    return shape

def add_textbox(slide, left, top, width, height, text, font_size,
                bold=False, italic=False, color=WHITE, align=PP_ALIGN.LEFT,
                font_name="Calibri", wrap=True):
    txb = slide.shapes.add_textbox(left, top, width, height)
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(font_size)
    run.font.bold  = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name  = font_name
    return txb

def add_para(tf, text, font_size, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, space_before=0, indent_level=0,
             font_name="Calibri", italic=False):
    from pptx.util import Pt
    p = tf.add_paragraph()
    p.alignment = align
    p.level = indent_level
    if space_before:
        p.space_before = Pt(space_before)
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(font_size)
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.name  = font_name
    run.font.italic = italic
    return p

def add_logo_icon(slide, left, top, size=Inches(0.55)):
    if os.path.exists(ICON_PATH):
        slide.shapes.add_picture(ICON_PATH, left, top, size, size)

def footer_bar(slide, text="BPMN IQ 2  |  Process Intelligence Platform  |  AT&T Enterprise Architecture"):
    add_rect(slide, 0, H - Inches(0.35), W, Inches(0.35), NAVY)
    add_textbox(slide, Inches(0.2), H - Inches(0.33), W - Inches(0.4), Inches(0.3),
                text, 7, color=LGREY, align=PP_ALIGN.LEFT)
    # Page number placeholder (right side)
    add_textbox(slide, W - Inches(1.2), H - Inches(0.33), Inches(1.1), Inches(0.3),
                "May 2026", 7, color=LGREY, align=PP_ALIGN.RIGHT)

def cyan_accent_line(slide, y_top=Inches(1.55)):
    add_rect(slide, 0, y_top, W, Inches(0.045), CYAN)

def slide_header(slide, title, subtitle=None):
    """Standard content slide header."""
    add_rect(slide, 0, 0, W, Inches(1.5), NAVY)
    add_logo_icon(slide, Inches(0.15), Inches(0.45))
    add_textbox(slide, Inches(0.85), Inches(0.18), W - Inches(1.1), Inches(0.7),
                title, 26, bold=True, color=WHITE)
    if subtitle:
        add_textbox(slide, Inches(0.85), Inches(0.84), W - Inches(1.1), Inches(0.5),
                    subtitle, 13, color=CYAN, italic=True)
    cyan_accent_line(slide, Inches(1.5))
    footer_bar(slide)

# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 1 – TITLE
# ══════════════════════════════════════════════════════════════════════════════
def slide_title():
    s = prs.slides.add_slide(blank_layout)

    # Full-bleed navy background
    add_rect(s, 0, 0, W, H, NAVY)

    # Left cyan accent stripe
    add_rect(s, 0, 0, Inches(0.18), H, CYAN)

    # Blue diagonal accent block (top-right)
    add_rect(s, W - Inches(4.5), 0, Inches(4.5), Inches(3.2), BLUE)

    # Logo banner
    if os.path.exists(LOGO_PATH):
        s.shapes.add_picture(LOGO_PATH, Inches(0.5), Inches(1.5), Inches(7), Inches(2.1))

    # Tagline
    add_textbox(s, Inches(0.5), Inches(3.75), Inches(8), Inches(0.5),
                "Process Intelligence Platform for Enterprise Architecture",
                18, italic=True, color=CYAN)

    # Divider rule
    add_rect(s, Inches(0.5), Inches(4.35), Inches(7.5), Inches(0.04), BLUE)

    # Sub-line
    add_textbox(s, Inches(0.5), Inches(4.5), Inches(8), Inches(0.4),
                "BPMN Authoring  ·  Application Risk Analytics  ·  Enterprise Cost Insights",
                13, color=LGREY)

    # Date / version
    add_textbox(s, Inches(0.5), H - Inches(0.8), Inches(5), Inches(0.4),
                "May 2026  |  v2.0", 12, color=MGREY)

    # Icon top-right
    add_logo_icon(s, W - Inches(2.8), Inches(0.3), Inches(2.4))

slide_title()

# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 2 – AGENDA / TABLE OF CONTENTS
# ══════════════════════════════════════════════════════════════════════════════
def slide_agenda():
    s = prs.slides.add_slide(blank_layout)
    add_rect(s, 0, 0, W, H, WHITE)
    slide_header(s, "Agenda", "What we'll cover today")

    items = [
        ("01", "What is BPMN IQ 2?",               "Purpose & vision"),
        ("02", "The Problem We Solve",              "EA pain points addressed"),
        ("03", "Feature Overview",                  "Capabilities & functions"),
        ("04", "Data Architecture",                 "Collections & data model"),
        ("05", "Technical Architecture",            "Tech stack & deployment"),
        ("06", "EA Framework Alignment",            "TOGAF, TM Forum, Capability Maps"),
        ("07", "Analytics Dashboard",               "Risk, compliance & cost insights"),
        ("08", "User Roles & Access Control",       "Security & governance model"),
        ("09", "Value Delivered",                   "Business outcomes"),
    ]

    col_w = Inches(6.4)
    for i, (num, title, sub) in enumerate(items):
        row = i % 5
        col = i // 5
        y = Inches(1.75) + row * Inches(1.02)
        x = Inches(0.4) + col * col_w

        add_rect(s, x, y, Inches(0.55), Inches(0.55), BLUE)
        add_textbox(s, x, y + Inches(0.05), Inches(0.55), Inches(0.45),
                    num, 14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_textbox(s, x + Inches(0.65), y, Inches(col_w - Inches(0.8)), Inches(0.35),
                    title, 13, bold=True, color=NAVY)
        add_textbox(s, x + Inches(0.65), y + Inches(0.32), Inches(col_w - Inches(0.8)), Inches(0.28),
                    sub, 10, color=MGREY)

slide_agenda()

# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 3 – WHAT IS BPMN IQ 2?
# ══════════════════════════════════════════════════════════════════════════════
def slide_what():
    s = prs.slides.add_slide(blank_layout)
    add_rect(s, 0, 0, W, H, WHITE)
    slide_header(s, "What is BPMN IQ 2?", "A purpose-built Enterprise Architecture intelligence platform")

    # Large pull-quote box
    add_rect(s, Inches(0.4), Inches(1.75), Inches(8.0), Inches(1.55), NAVY)
    add_rect(s, Inches(0.4), Inches(1.75), Inches(0.12), Inches(1.55), CYAN)
    add_textbox(s, Inches(0.65), Inches(1.85), Inches(7.6), Inches(1.3),
                "BPMN IQ 2 is a web-based platform that allows Enterprise Architects to author, store, "
                "and analyse BPMN 2.0 business process diagrams — and then derive structured intelligence "
                "about the applications, tasks, costs, and compliance risk embedded within those processes.",
                13, color=WHITE, wrap=True)

    # Three pillars
    pillars = [
        (BLUE,  "Author",  "Draw & save BPMN 2.0 diagrams with rich metadata tagging. "
                            "Classify by Line of Business, Channel, Domain, Product & Business Flow."),
        (CYAN,  "Enrich",  "Bind applications, actors, capabilities & cost data to every task "
                            "in every flow. Build a living map of your application landscape."),
        (NAVY,  "Analyse", "Dashboards surface risk scores, compliance exposure (CPNI, SPI, PCI, SOX), "
                            "application criticality & annual cost by flow or task."),
    ]
    pw = Inches(3.9)
    for i, (col, head, body) in enumerate(pillars):
        x = Inches(0.4) + i * Inches(4.25)
        add_rect(s, x, Inches(3.5), pw, Inches(0.5), col)
        add_textbox(s, x + Inches(0.15), Inches(3.52), pw - Inches(0.3), Inches(0.46),
                    head, 16, bold=True, color=WHITE)
        add_rect(s, x, Inches(4.0), pw, Inches(2.55), LGREY)
        add_textbox(s, x + Inches(0.15), Inches(4.1), pw - Inches(0.3), Inches(2.35),
                    body, 11, color=DGREY, wrap=True)

slide_what()

# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 4 – THE PROBLEM WE SOLVE
# ══════════════════════════════════════════════════════════════════════════════
def slide_problem():
    s = prs.slides.add_slide(blank_layout)
    add_rect(s, 0, 0, W, H, WHITE)
    slide_header(s, "The Problem We Solve", "Common Enterprise Architecture gaps BPMN IQ 2 addresses")

    problems = [
        ("Disconnected Process & App Data",
         "Process diagrams live in Visio or Confluence. Application data lives in LeanIX or ServiceNow. "
         "There is no tool that joins them — until now."),
        ("No Compliance Visibility at Process Level",
         "Architects cannot easily answer: 'Which business flows touch CPNI data?' or "
         "'What is the PCI exposure of the Checkout flow?' BPMN IQ 2 answers these instantly."),
        ("Opaque Application Costs",
         "Annual OpEx and DevEx are tracked in spreadsheets, never linked to the processes "
         "they support. BPMN IQ 2 surfaces cost-per-flow and cost-per-task."),
        ("Stale Capability Mapping",
         "Business capability models are rarely kept in sync with running processes. "
         "BPMN IQ 2 auto-matches tasks to capabilities using AI-assisted scoring."),
        ("No Lifecycle Awareness",
         "Architects cannot tell which flows depend on retiring or sunset applications. "
         "BPMN IQ 2 flags lifecycle risk per task and per flow."),
        ("Governance Without Guardrails",
         "Diagram quality and metadata completeness are inconsistent. "
         "BPMN IQ 2 enforces state workflows (Draft → Review → Published → Retired)."),
    ]

    for i, (title, body) in enumerate(problems):
        row = i % 3
        col = i // 3
        x = Inches(0.4) + col * Inches(6.45)
        y = Inches(1.75) + row * Inches(1.72)
        add_rect(s, x, y, Inches(6.1), Inches(1.55), RGBColor(0xF0, 0xF8, 0xFF))
        add_rect(s, x, y, Inches(0.1), Inches(1.55), BLUE)
        add_textbox(s, x + Inches(0.22), y + Inches(0.08), Inches(5.7), Inches(0.38),
                    title, 12, bold=True, color=NAVY)
        add_textbox(s, x + Inches(0.22), y + Inches(0.46), Inches(5.7), Inches(1.0),
                    body, 10, color=DGREY, wrap=True)

slide_problem()

# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 5 – FEATURE OVERVIEW
# ══════════════════════════════════════════════════════════════════════════════
def slide_features():
    s = prs.slides.add_slide(blank_layout)
    add_rect(s, 0, 0, W, H, WHITE)
    slide_header(s, "Feature Overview", "What BPMN IQ 2 can do")

    features = [
        ("BPMN 2.0 Canvas Editor",
         "Full bpmn-js powered diagram editor with lanes, gateways, events, sub-processes "
         "and custom extension attributes. Auto-saves to MongoDB."),
        ("Business Flow Tagging",
         "Each diagram is classified by Line of Business, Channel, Domain, Subdomain, "
         "Product, and Business Flow — enabling multi-dimensional filtering."),
        ("Application Binding",
         "Tasks are linked to enterprise applications. Every application carries lifecycle status, "
         "criticality, type, cost data, and compliance flags."),
        ("Capability Auto-Match",
         "An AI-assisted engine matches diagram tasks to the Business Capability Model, "
         "returning a confidence score and justification for each match."),
        ("Analytics Dashboard",
         "Interactive charts: Top-20 cost bars (flow & task), Application Criticality "
         "pie, Compliance Radar, 3D Flow Explorer, and Risk Score tables."),
        ("Compliance Risk Scoring",
         "Each task receives a composite risk score based on CPNI, SPI, PCI, SOX, "
         "customer-facing, and internet-facing flags across its applications."),
        ("Cost Intelligence",
         "Annual operational and development cost per flow and per task, sourced from "
         "the application reference data and visualised for 2016-2025."),
        ("Factory Editors",
         "Dedicated CRUD editors for Applications, Business Flows, Tasks, Actors, "
         "Capabilities, Products, Lines of Business, Channels, Domains, and Subdomains."),
        ("Role-Based Access Control",
         "Users are assigned Roles (Admin, Editor, Viewer). Each role has a Capability-Permission "
         "matrix. Session-managed via server-side tokens with auto-expiry."),
        ("State Workflow",
         "Diagrams and assets move through a governed lifecycle: "
         "Draft → In Review → Published → Deprecated → Retired."),
        ("Detailed Cost Reports",
         "Generate HTML reports showing per-process application cost breakdown with "
         "operational vs. development split for any business flow."),
        ("Search & Discover",
         "Full-text search across diagram names, tags, and content. "
         "Filter by any taxonomy dimension across the library."),
    ]

    col_w = Inches(4.2)
    for i, (title, body) in enumerate(features):
        row = i % 4
        col = i // 4
        x = Inches(0.3) + col * Inches(4.35)
        y = Inches(1.72) + row * Inches(1.38)
        add_rect(s, x, y, col_w, Inches(1.28), RGBColor(0xF4, 0xF9, 0xFF))
        add_rect(s, x, y, col_w, Inches(0.08), BLUE)
        add_textbox(s, x + Inches(0.1), y + Inches(0.12), col_w - Inches(0.2), Inches(0.35),
                    title, 11, bold=True, color=NAVY)
        add_textbox(s, x + Inches(0.1), y + Inches(0.46), col_w - Inches(0.2), Inches(0.75),
                    body, 9, color=DGREY, wrap=True)

slide_features()

# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 6 – DATA ARCHITECTURE
# ══════════════════════════════════════════════════════════════════════════════
def slide_data():
    s = prs.slides.add_slide(blank_layout)
    add_rect(s, 0, 0, W, H, WHITE)
    slide_header(s, "Data Architecture", "MongoDB collections & key fields")

    collections = [
        ("diagrams",        NAVY,  ["name, description, tags, version",
                                    "xml (BPMN 2.0 source)",
                                    "capabilities[], changeHistory[]",
                                    "tasks[].applications[].annualCosts[]"]),
        ("tasks",           BLUE,  ["name, businessFlow, product",
                                    "domain, subdomain, channel, actor",
                                    "applications[] (names)",
                                    "sequence, owner, state"]),
        ("applications",    CYAN,  ["name, lifecycleStatus, applicationType",
                                    "businessCriticality, customerFacing",
                                    "cpniIndicator, handleSpi, pciData, soxFsa",
                                    "annualTotalCosts[0..9] (2016–2025)"]),
        ("businessflows",   NAVY,  ["name, lineOfBusiness, channel",
                                    "domain, subdomain, product",
                                    "tasks[] with applications[].annualCosts",
                                    "Aggregated at import time"]),
        ("capabilities",    BLUE,  ["capabilityId, name, domainName, aspect",
                                    "briefDescription, fullDescription",
                                    "tmfStatus, tmfVersion",
                                    "owner, state"]),
        ("actors",          CYAN,  ["name, role, description", "owner, state"]),
        ("roles",           NAVY,  ["name, capabilities[]{function, permission}",
                                    "Defines RBAC matrix"]),
        ("users",           BLUE,  ["username, passwordHash",
                                    "role, lastLogin"]),
        ("states",          CYAN,  ["Workflow state definitions",
                                    "Draft → InReview → Published → Retired"]),
    ]

    cw = Inches(1.4)
    ch = Inches(1.6)
    pad = Inches(0.08)
    cols_per_row = 5
    for i, (name, color, fields) in enumerate(collections):
        row = i // cols_per_row
        col = i % cols_per_row
        x = Inches(0.3) + col * (cw + pad)
        y = Inches(1.72) + row * (ch + Inches(0.22))
        add_rect(s, x, y, cw, Inches(0.42), color)
        add_textbox(s, x + Inches(0.07), y + Inches(0.06), cw - Inches(0.14), Inches(0.34),
                    name, 11, bold=True, color=WHITE)
        add_rect(s, x, y + Inches(0.42), cw, ch - Inches(0.42), RGBColor(0xF2, 0xF8, 0xFD))
        for j, field in enumerate(fields):
            add_textbox(s, x + Inches(0.08), y + Inches(0.46) + j * Inches(0.26),
                        cw - Inches(0.12), Inches(0.26),
                        "• " + field, 8, color=DGREY)

    # Note
    add_textbox(s, Inches(0.3), H - Inches(0.9), Inches(12), Inches(0.35),
                "All collections reside in MongoDB (bpmn_iq database).  "
                "businessflows is a denormalised read-model rebuilt on diagram import/save.",
                9, color=MGREY, italic=True)

slide_data()

# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 7 – TECHNICAL ARCHITECTURE
# ══════════════════════════════════════════════════════════════════════════════
def slide_tech():
    s = prs.slides.add_slide(blank_layout)
    add_rect(s, 0, 0, W, H, WHITE)
    slide_header(s, "Technical Architecture", "Full-stack technology choices")

    layers = [
        ("Presentation Layer",  NAVY,  Inches(1.72),
         [("React 18 + TypeScript",   "Component-based SPA. Vite dev server (port 5173)."),
          ("Ant Design 5",            "UI kit: Tabs, Select, Card, Table, Statistic, Modal."),
          ("bpmn-js",                 "Industry-standard BPMN 2.0 renderer & modeller."),
          ("Recharts",                "BarChart, PieChart, RadarChart — all responsive."),
          ("react-plotly.js",         "3D scatter / surface charts for Flow Explorer."),]),
        ("API / Business Layer", BLUE,  Inches(3.28),
         [("Node.js 20 + Express",    "RESTful API server (port 3001). JWT-like session tokens."),
          ("Mongoose ODM",            "Schema validation, indexes, text search."),
          ("Route modules",           "diagrams, tasks, dashboard, reports, auth, admin, capabilities, actors."),
          ("AI Capability Match",     "Embeddings-style scoring against Business Capability Model."),]),
        ("Data Layer",          CYAN,  Inches(4.84),
         [("MongoDB 7",               "Document store – bpmn_iq database, 15 collections."),
          ("Application Ref Data",    "5 000+ app records sourced from app_ref_data.csv."),
          ("BPMN Extension Schema",   "bpmniq-moddle.json extends BPMN 2.0 with custom attributes."),]),
    ]

    layer_h = Inches(1.35)
    for (label, color, y, items) in layers:
        add_rect(s, Inches(0.3), y, Inches(2.1), layer_h, color)
        add_textbox(s, Inches(0.34), y + Inches(0.45), Inches(2.0), Inches(0.55),
                    label, 12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_rect(s, Inches(2.4), y, Inches(10.5), layer_h, RGBColor(0xF3, 0xF8, 0xFD))
        for j, (tech, desc) in enumerate(items):
            tx = Inches(2.55) + j * Inches(2.08)
            add_rect(s, tx, y + Inches(0.12), Inches(1.9), Inches(1.12), WHITE)
            add_rect(s, tx, y + Inches(0.12), Inches(1.9), Inches(0.04), color)
            add_textbox(s, tx + Inches(0.08), y + Inches(0.18), Inches(1.75), Inches(0.32),
                        tech, 9, bold=True, color=NAVY)
            add_textbox(s, tx + Inches(0.08), y + Inches(0.48), Inches(1.75), Inches(0.7),
                        desc, 8, color=DGREY, wrap=True)

    # Ports legend
    add_textbox(s, Inches(0.3), H - Inches(0.9), Inches(12), Inches(0.32),
                "Ports:  Client → :5173 (Vite)   API → :3001 (Express)   DB → :27017 (MongoDB)   "
                "Proxy: Vite dev-server forwards /api/* → Express",
                9, color=MGREY, italic=True)

slide_tech()

# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 8 – EA FRAMEWORK ALIGNMENT
# ══════════════════════════════════════════════════════════════════════════════
def slide_ea():
    s = prs.slides.add_slide(blank_layout)
    add_rect(s, 0, 0, W, H, WHITE)
    slide_header(s, "Enterprise Architecture Alignment",
                 "How BPMN IQ 2 maps to industry EA frameworks")

    frameworks = [
        ("TOGAF ADM", NAVY,
         ["Business Architecture phase — process diagrams\n  directly populate the BA artefact set.",
          "Information Systems Architecture — applications\n  linked to every task form the IS layer.",
          "Technology Architecture — lifecycle & cost data\n  feeds Technology Roadmap decisions.",
          "Architecture Governance — state workflow\n  (Draft→Published→Retired) enforces ADM gates."]),
        ("TM Forum eTOM / SID", BLUE,
         ["Processes classified by Domain & Subdomain\n  align to eTOM Level 2/3 groupings.",
          "Applications mapped to eTOM process areas\n  enable TM Forum Frameworx compliance.",
          "Capability IDs carry tmfStatus & tmfVersion\n  for direct eTOM traceability.",
          "Business Flow hierarchy mirrors\n  eTOM Engaged Party & Customer domains."]),
        ("Business Capability Model", CYAN,
         ["AI-assisted capability matching scores\n  each BPMN task against the BCM.",
          "Confidence scores & justifications\n  stored per diagram for audit.",
          "Capability Factory enables BCM maintenance\n  inside the same platform.",
          "Supports heat-mapping: which capabilities\n  have the most regulatory exposure?"]),
        ("Application Portfolio Mgmt", NAVY,
         ["Every application carries LeanIX-compatible\n  lifecycle, criticality & type metadata.",
          "Cost per application per year (2016-2025)\n  supports TCO & run-cost analysis.",
          "Process-to-application dependency matrix\n  derived automatically from task bindings.",
          "Retiring-application risk surfaced in\n  Dashboard compliance radar."]),
    ]

    fw = Inches(3.15)
    for i, (name, color, points) in enumerate(frameworks):
        col = i % 2
        row = i // 2
        x = Inches(0.3) + col * Inches(6.45)
        y = Inches(1.72) + row * Inches(2.45)
        add_rect(s, x, y, fw * 2 + Inches(0.1), Inches(0.44), color)
        add_textbox(s, x + Inches(0.12), y + Inches(0.06), fw * 2, Inches(0.36),
                    name, 14, bold=True, color=WHITE)
        for j, pt in enumerate(points):
            py = y + Inches(0.52) + j * Inches(0.44)
            add_rect(s, x + Inches(0.1), py + Inches(0.1), Inches(0.12), Inches(0.12), color)
            add_textbox(s, x + Inches(0.32), py, fw * 2 - Inches(0.42), Inches(0.44),
                        pt, 9, color=DGREY, wrap=True)

slide_ea()

# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 9 – ANALYTICS DASHBOARD
# ══════════════════════════════════════════════════════════════════════════════
def slide_dashboard():
    s = prs.slides.add_slide(blank_layout)
    add_rect(s, 0, 0, W, H, WHITE)
    slide_header(s, "Analytics Dashboard", "Insights derived from the process & application data")

    sections = [
        ("Business Flow Comparison", NAVY,
         [("Top 20 Flows by Cost (2025)",
           "Stacked bar: Op Cost vs Dev Cost per flow. Sorted by total. Hover for exact values."),
          ("Compliance Bar Chart",
           "App compliance counts (CPNI, SPI, PCI, SOX, Customer/Internet-facing) stacked per flow."),
          ("Application Criticality Pie",
           "Distribution of Mission Critical / Business Critical / Standard apps. "
           "Filter by one or more flows using the multi-select picker."),
          ("Compliance Radar",
           "Spider chart of compliance dimensions for up to 5 selected flows. "
           "Defaults to Top 5 by risk score.")]),
        ("Task Comparison", BLUE,
         [("Top 20 Tasks by Cost (2025)",
           "Same stacked bar format as flows. Shows task name + owning business flow in tooltip."),
          ("Task Compliance Bar",
           "Per-task compliance flag counts across the riskiest tasks."),
          ("Application Criticality Pie",
           "Filtered by task selection (multi-select, 710 unique task names). "
           "Aggregates all instances of a shared task name across flows."),
          ("Task Compliance Radar",
           "Up to 5 selected tasks. Defaults to Top 5 riskiest. Dynamic title on selection.")]),
        ("3D Flow Explorer", CYAN,
         [("3D Scatter / Surface",
           "Business flows plotted in 3D space: X=app count, Y=risk score, Z=cost. "
           "Powered by Plotly.js. Rotate, zoom, hover for detail."),
          ("Summary Statistics",
           "KPI cards: Total Flows, Total Tasks, Total Apps, Avg Risk Score, "
           "High Risk Count, Max Risk Score.")]),
    ]

    sy = Inches(1.72)
    for sec_i, (sec_name, color, items) in enumerate(sections):
        x = Inches(0.3) + sec_i * Inches(4.35)
        sw = Inches(4.1)
        add_rect(s, x, sy, sw, Inches(0.4), color)
        add_textbox(s, x + Inches(0.1), sy + Inches(0.05), sw - Inches(0.2), Inches(0.32),
                    sec_name, 12, bold=True, color=WHITE)
        for j, (chart, desc) in enumerate(items):
            cy2 = sy + Inches(0.48) + j * Inches(1.18)
            add_rect(s, x, cy2, sw, Inches(1.08), RGBColor(0xF0, 0xF7, 0xFF) if sec_i % 2 == 0 else RGBColor(0xF5, 0xF5, 0xFF))
            add_rect(s, x, cy2, Inches(0.06), Inches(1.08), color)
            add_textbox(s, x + Inches(0.14), cy2 + Inches(0.06), sw - Inches(0.2), Inches(0.3),
                        chart, 10, bold=True, color=NAVY)
            add_textbox(s, x + Inches(0.14), cy2 + Inches(0.35), sw - Inches(0.2), Inches(0.68),
                        desc, 9, color=DGREY, wrap=True)

slide_dashboard()

# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 10 – USER ROLES & ACCESS CONTROL
# ══════════════════════════════════════════════════════════════════════════════
def slide_roles():
    s = prs.slides.add_slide(blank_layout)
    add_rect(s, 0, 0, W, H, WHITE)
    slide_header(s, "User Roles & Access Control", "Security, governance & workflow model")

    # Role cards
    roles = [
        ("Admin",   NAVY,
         ["Full platform access",
          "Manage users, roles & permissions",
          "Access all factory editors",
          "Publish / retire any diagram",
          "Reset passwords"]),
        ("Editor",  BLUE,
         ["Author & edit diagrams",
          "Bind applications & capabilities",
          "Submit diagrams for review",
          "Edit own factory records",
          "View analytics dashboard"]),
        ("Viewer",  CYAN,
         ["Read-only diagram access",
          "View factory reference data",
          "View analytics & reports",
          "No create / edit / delete",
          "No admin functions"]),
    ]

    rw = Inches(4.0)
    for i, (role, color, perms) in enumerate(roles):
        x = Inches(0.3) + i * Inches(4.35)
        add_rect(s, x, Inches(1.75), rw, Inches(0.55), color)
        add_textbox(s, x + Inches(0.12), Inches(1.8), rw - Inches(0.2), Inches(0.45),
                    role, 18, bold=True, color=WHITE)
        for j, perm in enumerate(perms):
            py = Inches(2.42) + j * Inches(0.42)
            add_rect(s, x + Inches(0.15), py + Inches(0.12), Inches(0.14), Inches(0.14), color)
            add_textbox(s, x + Inches(0.4), py, rw - Inches(0.5), Inches(0.38),
                        perm, 11, color=DGREY)

    # State workflow section
    add_rect(s, Inches(0.3), Inches(4.8), Inches(12.7), Inches(0.38), NAVY)
    add_textbox(s, Inches(0.5), Inches(4.86), Inches(12), Inches(0.28),
                "ASSET STATE WORKFLOW", 11, bold=True, color=WHITE)

    states = [("Draft", MGREY), ("In Review", BLUE), ("Published", CYAN),
              ("Deprecated", ACCENT), ("Retired", RGBColor(0xCC,0x33,0x33))]
    sw_box = Inches(2.3)
    for i, (state, color) in enumerate(states):
        sx = Inches(0.4) + i * (sw_box + Inches(0.1))
        add_rect(s, sx, Inches(5.28), sw_box, Inches(0.46), color)
        add_textbox(s, sx, Inches(5.32), sw_box, Inches(0.38),
                    state, 12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        if i < 4:
            add_textbox(s, sx + sw_box + Inches(0.0), Inches(5.36), Inches(0.12), Inches(0.3),
                        "→", 14, bold=True, color=NAVY)

    add_textbox(s, Inches(0.4), Inches(5.88), Inches(12), Inches(0.35),
                "Applies to: Diagrams · Tasks · Applications · Capabilities · Actors · Business Flows",
                10, color=MGREY, italic=True)

slide_roles()

# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 11 – VALUE DELIVERED
# ══════════════════════════════════════════════════════════════════════════════
def slide_value():
    s = prs.slides.add_slide(blank_layout)
    add_rect(s, 0, 0, W, H, NAVY)
    add_rect(s, 0, 0, Inches(0.18), H, CYAN)

    add_textbox(s, Inches(0.4), Inches(0.3), Inches(12), Inches(0.65),
                "Value Delivered", 32, bold=True, color=WHITE)
    add_textbox(s, Inches(0.4), Inches(0.95), Inches(12), Inches(0.38),
                "What BPMN IQ 2 enables for the Enterprise Architecture practice",
                14, italic=True, color=CYAN)
    add_rect(s, Inches(0.4), Inches(1.4), Inches(12.5), Inches(0.04), BLUE)

    values = [
        ("Single Source of Truth",
         "One authoritative store for all business process diagrams, "
         "their applications, tasks, costs, and compliance metadata."),
        ("Faster Compliance Answers",
         "Identify all flows touching CPNI / PCI / SOX data in seconds, "
         "not days of spreadsheet analysis."),
        ("Informed Investment Decisions",
         "See the true cost of each business flow before approving "
         "modernisation, migration, or retirement projects."),
        ("Living Capability Map",
         "Capabilities stay in sync with running processes — no manual "
         "reconciliation between BCM and the process library."),
        ("Reduced Architecture Debt",
         "State workflow prevents unreviewed or outdated diagrams "
         "from polluting the reference architecture."),
        ("EA-Led Risk Management",
         "Risk scores flag which flows carry the highest combination of "
         "regulated data, sunset apps, and internet exposure."),
    ]

    vw = Inches(4.1)
    for i, (title, body) in enumerate(values):
        row = i % 3
        col = i // 3
        x = Inches(0.4) + col * Inches(6.45)
        y = Inches(1.65) + row * Inches(1.75)
        add_rect(s, x, y, vw * 1.55, Inches(1.55), RGBColor(0x00, 0x3A, 0x7A))
        add_rect(s, x, y, Inches(0.1), Inches(1.55), CYAN)
        add_textbox(s, x + Inches(0.22), y + Inches(0.1), vw * 1.5 - Inches(0.25), Inches(0.38),
                    title, 13, bold=True, color=CYAN)
        add_textbox(s, x + Inches(0.22), y + Inches(0.48), vw * 1.5 - Inches(0.25), Inches(1.0),
                    body, 10, color=LGREY, wrap=True)

    footer_bar(s)

slide_value()

# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 12 – THANK YOU / CLOSE
# ══════════════════════════════════════════════════════════════════════════════
def slide_close():
    s = prs.slides.add_slide(blank_layout)
    add_rect(s, 0, 0, W, H, NAVY)
    add_rect(s, 0, 0, Inches(0.18), H, CYAN)
    add_rect(s, W - Inches(5), 0, Inches(5), H, BLUE)

    if os.path.exists(LOGO_PATH):
        s.shapes.add_picture(LOGO_PATH, Inches(0.5), Inches(1.8), Inches(7.2), Inches(2.15))

    add_textbox(s, Inches(0.5), Inches(4.15), Inches(7), Inches(0.5),
                "Process Intelligence Platform", 18, italic=True, color=CYAN)

    add_rect(s, Inches(0.5), Inches(4.75), Inches(6.5), Inches(0.04), BLUE)

    add_textbox(s, Inches(0.5), Inches(4.92), Inches(6.5), Inches(0.4),
                "Questions?  Reach out to the EA Tools team.",
                14, color=LGREY)

    add_textbox(s, Inches(0.5), Inches(5.45), Inches(6.5), Inches(0.35),
                "http://localhost:5173  (development)  |  May 2026",
                11, color=MGREY)

    add_logo_icon(s, W - Inches(3.5), Inches(1.8), Inches(3.0))

    footer_bar(s)

slide_close()

# ══════════════════════════════════════════════════════════════════════════════
#  SAVE
# ══════════════════════════════════════════════════════════════════════════════
prs.save(OUT_PATH)
print(f"Saved: {OUT_PATH}")
print(f"Slides: {len(prs.slides)}")
